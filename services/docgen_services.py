"""
services/docgen_services.py
Complete document generation service supporting PDF, DOCX, XLSX, and PPTX formats.
"""

import os
from io import BytesIO
from typing import Dict, List, Any, Optional, Union

# Third-party imports
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle, Image as RLImage
from reportlab.lib import colors
from reportlab.pdfgen import canvas

from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from pptx import Presentation
from pptx.util import Inches as PPTXInches, Pt as PPTPt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor as PPTXRGBColor

from utils.file_utils import (
    ensure_generated_dir,
    generate_unique_filename,
    get_full_document_path,
    decode_base64_image,
)


# ============================================================================
# PDF Generator
# ============================================================================

def _generate_pdf(title: str, content: List[Dict[str, Any]]) -> str:
    """
    Generate a PDF document using ReportLab Platypus.
    
    Args:
        title: Document title
        content: List of content blocks
    
    Returns:
        Absolute path to generated PDF
    """
    filename = generate_unique_filename("pdf")
    filepath = get_full_document_path(filename)
    
    # Create PDF
    doc = SimpleDocTemplate(
        filepath,
        pagesize=letter,
        rightMargin=1 * inch,
        leftMargin=1 * inch,
        topMargin=1 * inch,
        bottomMargin=1 * inch,
    )
    
    # Get stylesheet
    styles = getSampleStyleSheet()
    
    # Custom styles
    custom_heading1 = ParagraphStyle(
        'CustomHeading1',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#1a1a1a'),
        spaceAfter=12,
        spaceBefore=6,
    )
    
    custom_heading2 = ParagraphStyle(
        'CustomHeading2',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=colors.HexColor('#333333'),
        spaceAfter=10,
        spaceBefore=6,
    )
    
    custom_body = ParagraphStyle(
        'CustomBody',
        parent=styles['BodyText'],
        fontSize=11,
        spaceAfter=8,
        spaceBefore=6,
    )
    
    # Story for PDF elements
    story = []
    
    # Add title
    story.append(Paragraph(title, custom_heading1))
    story.append(Spacer(1, 0.2 * inch))
    
    # Process content blocks
    for block in content:
        if "h1" in block:
            story.append(Paragraph(block["h1"], custom_heading1))
            story.append(Spacer(1, 0.1 * inch))
        
        elif "h2" in block:
            story.append(Paragraph(block["h2"], custom_heading2))
            story.append(Spacer(1, 0.1 * inch))
        
        elif "p" in block:
            story.append(Paragraph(block["p"], custom_body))
            story.append(Spacer(1, 0.1 * inch))
        
        elif "bullet" in block:
            bullets = block["bullet"]
            for bullet_text in bullets:
                story.append(Paragraph(f"• {bullet_text}", custom_body))
            story.append(Spacer(1, 0.1 * inch))
        
        elif "ul" in block:
            # Handle unordered list (ul) same as bullet
            items = block["ul"]
            for item_text in items:
                story.append(Paragraph(f"• {item_text}", custom_body))
            story.append(Spacer(1, 0.1 * inch))
        
        elif "table" in block:
            table_data = block["table"]
            if table_data:
                # Convert to Platypus Table
                table = Table(table_data, colWidths=[2.5 * inch] * len(table_data[0]))
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4a4a4a')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('FONTSIZE', (0, 1), (-1, -1), 10),
                ]))
                story.append(table)
                story.append(Spacer(1, 0.2 * inch))
        
        elif "image" in block:
            # Handle base64 or file path images
            image_data = block["image"]
            if image_data.startswith("data:"):
                # Base64 image
                image_path = decode_base64_image(image_data)
            else:
                # Assume it's a file path
                image_path = image_data
            
            if os.path.exists(image_path):
                try:
                    img = RLImage(image_path, width=4 * inch, height=3 * inch)
                    story.append(img)
                    story.append(Spacer(1, 0.2 * inch))
                except Exception:
                    pass
    
    # Build PDF
    doc.build(story)
    return filepath


# ============================================================================
# DOCX Generator
# ============================================================================

def _generate_docx(title: str, content: List[Dict[str, Any]]) -> str:
    """
    Generate a DOCX document using python-docx.
    
    Args:
        title: Document title
        content: List of content blocks
    
    Returns:
        Absolute path to generated DOCX
    """
    filename = generate_unique_filename("docx")
    filepath = get_full_document_path(filename)
    
    doc = Document()
    
    # Set document margins
    sections = doc.sections
    for section in sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1)
        section.right_margin = Inches(1)
    
    # Add title
    title_para = doc.add_paragraph(title, style='Heading 1')
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Process content blocks
    for block in content:
        if "h1" in block:
            para = doc.add_paragraph(block["h1"], style='Heading 1')
            para.paragraph_format.space_after = Pt(12)
            para.paragraph_format.space_before = Pt(6)
        
        elif "h2" in block:
            para = doc.add_paragraph(block["h2"], style='Heading 2')
            para.paragraph_format.space_after = Pt(10)
            para.paragraph_format.space_before = Pt(6)
        
        elif "p" in block:
            para = doc.add_paragraph(block["p"], style='Normal')
            para.paragraph_format.space_after = Pt(8)
            para.paragraph_format.space_before = Pt(6)
        
        elif "bullet" in block:
            bullets = block["bullet"]
            for bullet_text in bullets:
                para = doc.add_paragraph(bullet_text, style='List Bullet')
                para.paragraph_format.space_after = Pt(4)
        
        elif "ul" in block:
            # Handle unordered list (ul) same as bullet
            items = block["ul"]
            for item_text in items:
                para = doc.add_paragraph(item_text, style='List Bullet')
                para.paragraph_format.space_after = Pt(4)
        
        elif "table" in block:
            table_data = block["table"]
            if table_data:
                rows = len(table_data)
                cols = len(table_data[0]) if rows > 0 else 0
                table = doc.add_table(rows=rows, cols=cols)
                table.style = 'Light Grid Accent 1'
                
                # Fill table
                for i, row in enumerate(table_data):
                    for j, cell_text in enumerate(row):
                        cell = table.rows[i].cells[j]
                        cell.text = str(cell_text)
                        
                        # Style header row
                        if i == 0:
                            for paragraph in cell.paragraphs:
                                for run in paragraph.runs:
                                    run.font.bold = True
                                    run.font.size = Pt(11)
                        
                        # Center alignment
                        for paragraph in cell.paragraphs:
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        elif "image" in block:
            # Handle base64 or file path images
            image_data = block["image"]
            if image_data.startswith("data:"):
                image_path = decode_base64_image(image_data)
            else:
                image_path = image_data
            
            if os.path.exists(image_path):
                try:
                    doc.add_picture(image_path, width=Inches(4))
                except Exception:
                    pass
    
    doc.save(filepath)
    return filepath


# ============================================================================
# XLSX Generator
# ============================================================================

def _generate_xlsx(title: str, content: List[Dict[str, Any]]) -> str:
    """
    Generate an XLSX document using openpyxl.
    
    Args:
        title: Document title
        content: List of content blocks
    
    Returns:
        Absolute path to generated XLSX
    """
    filename = generate_unique_filename("xlsx")
    filepath = get_full_document_path(filename)
    
    wb = Workbook()
    ws = wb.active
    ws.title = "Document"
    
    # Add title
    title_cell = ws['A1']
    title_cell.value = title
    title_cell.font = Font(name='Calibri', size=16, bold=True, color='FFFFFF')
    title_cell.fill = PatternFill(start_color='1a1a1a', end_color='1a1a1a', fill_type='solid')
    title_cell.alignment = Alignment(horizontal='center', vertical='center')
    ws.merge_cells('A1:D1')
    ws.row_dimensions[1].height = 30
    
    current_row = 3
    
    # Define styles
    thin_border = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin')
    )
    
    header_fill = PatternFill(start_color='D3D3D3', end_color='D3D3D3', fill_type='solid')
    
    # Process content blocks
    for block in content:
        if "h1" in block:
            ws[f'A{current_row}'].value = block["h1"]
            ws[f'A{current_row}'].font = Font(bold=True, size=14)
            ws.merge_cells(f'A{current_row}:D{current_row}')
            current_row += 1
        
        elif "h2" in block:
            ws[f'A{current_row}'].value = block["h2"]
            ws[f'A{current_row}'].font = Font(bold=True, size=12)
            ws.merge_cells(f'A{current_row}:D{current_row}')
            current_row += 1
        
        elif "p" in block:
            ws[f'A{current_row}'].value = block["p"]
            ws.merge_cells(f'A{current_row}:D{current_row}')
            current_row += 2
        
        elif "bullet" in block:
            bullets = block["bullet"]
            for bullet_text in bullets:
                ws[f'A{current_row}'].value = f"• {bullet_text}"
                current_row += 1
            current_row += 1
        
        elif "ul" in block:
            # Handle unordered list (ul) same as bullet
            items = block["ul"]
            for item_text in items:
                ws[f'A{current_row}'].value = f"• {item_text}"
                current_row += 1
            current_row += 1
        
        elif "table" in block:
            table_data = block["table"]
            if table_data:
                start_row = current_row
                for i, row in enumerate(table_data):
                    for j, cell_text in enumerate(row):
                        cell = ws.cell(row=current_row, column=j + 1)
                        cell.value = str(cell_text)
                        cell.border = thin_border
                        cell.alignment = Alignment(horizontal='center', vertical='center')
                        
                        # Style header row
                        if i == 0:
                            cell.font = Font(bold=True, color='FFFFFF')
                            cell.fill = PatternFill(start_color='4a4a4a', end_color='4a4a4a', fill_type='solid')
                        # Alternating row colors
                        elif i % 2 == 0:
                            cell.fill = PatternFill(start_color='F2F2F2', end_color='F2F2F2', fill_type='solid')
                    
                    current_row += 1
                current_row += 1
    
    # Auto-adjust column widths
    for col in ws.columns:
        max_length = 0
        column_letter = get_column_letter(col[0].column)
        for cell in col:
            try:
                max_length = max(max_length, len(str(cell.value)))
            except:
                pass
        ws.column_dimensions[column_letter].width = min(max_length + 2, 50)
    
    wb.save(filepath)
    return filepath


# ============================================================================
# PPTX Generator
# ============================================================================

def _generate_pptx(title: str, content: List[Dict[str, Any]]) -> str:
    """
    Generate a PPTX presentation using python-pptx.
    
    Args:
        title: Document title
        content: List of content blocks
    
    Returns:
        Absolute path to generated PPTX
    """
    filename = generate_unique_filename("pptx")
    filepath = get_full_document_path(filename)
    
    prs = Presentation()
    prs.slide_width = PPTXInches(10)
    prs.slide_height = PPTXInches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title
    subtitle.text = "Document Generated by DocGen"
    
    # Content slides
    content_text = ""
    for block in content:
        if "h1" in block:
            # Create new slide for h1
            if content_text.strip():
                blank_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(blank_layout)
                text_box = slide.shapes.add_textbox(
                    PPTXInches(0.5), PPTXInches(0.5),
                    PPTXInches(9), PPTXInches(6.5)
                )
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                p = text_frame.paragraphs[0]
                p.text = content_text
                p.font.size = PPTPt(12)
                content_text = ""
            
            # Add h1 as slide title
            blank_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(blank_layout)
            title_box = slide.shapes.add_textbox(
                PPTXInches(0.5), PPTXInches(0.5),
                PPTXInches(9), PPTXInches(1)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = block["h1"]
            p.font.size = PPTPt(32)
            p.font.bold = True
            p.font.color.rgb = PPTXRGBColor(26, 26, 26)
        
        elif "h2" in block:
            content_text += f"\n{block['h2']}\n"
        
        elif "p" in block:
            content_text += f"{block['p']}\n"
        
        elif "bullet" in block:
            bullets = block["bullet"]
            for bullet_text in bullets:
                content_text += f"• {bullet_text}\n"
        
        elif "ul" in block:
            # Handle unordered list (ul) same as bullet
            items = block["ul"]
            for item_text in items:
                content_text += f"• {item_text}\n"
        
        elif "table" in block:
            table_data = block["table"]
            if table_data:
                blank_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(blank_layout)
                
                rows = len(table_data)
                cols = len(table_data[0]) if rows > 0 else 0
                
                left = PPTXInches(0.5)
                top = PPTXInches(0.5)
                width = PPTXInches(9)
                height = PPTXInches(6)
                
                table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # Fill table
                for i, row in enumerate(table_data):
                    for j, cell_text in enumerate(row):
                        cell = table_shape.cell(i, j)
                        cell.text = str(cell_text)
                        
                        # Style header row
                        if i == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = PPTXRGBColor(74, 74, 74)
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.color.rgb = PPTXRGBColor(255, 255, 255)
                                    run.font.bold = True
    
    # Add remaining content if any
    if content_text.strip():
        blank_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(blank_layout)
        text_box = slide.shapes.add_textbox(
            PPTXInches(0.5), PPTXInches(0.5),
            PPTXInches(9), PPTXInches(6.5)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = content_text
        p.font.size = PPTPt(12)
    
    prs.save(filepath)
    return filepath


# ============================================================================
# Main Entry Function
# ============================================================================

def generate_document(doc_payload: Dict[str, Any]) -> str:
    """
    Main entry point for document generation.
    
    Accepts a JSON payload with:
    {
        "type": "pdf" | "docx" | "xlsx" | "pptx",
        "title": "string",
        "content": [
            {"h1": "Title"},
            {"h2": "Subtitle"},
            {"p": "Paragraph text"},
            {"table": [["A", "B"], ["1", "2"]]},
            {"bullet": ["item1", "item2"]},
            {"image": "base64_or_file_path"}
        ]
    }
    
    Returns:
        Absolute file path of the generated document.
    
    Raises:
        ValueError: If document type is not supported or payload is invalid.
    """
    # Validate payload
    if not isinstance(doc_payload, dict):
        raise ValueError("Payload must be a dictionary")
    
    doc_type = doc_payload.get("type", "").lower()
    title = doc_payload.get("title", "Untitled Document")
    content = doc_payload.get("content", [])
    
    if not isinstance(content, list):
        raise ValueError("Content must be a list of content blocks")
    
    # Ensure /generated directory exists
    ensure_generated_dir()
    
    # Generate document based on type
    if doc_type == "pdf":
        return _generate_pdf(title, content)
    
    elif doc_type == "docx":
        return _generate_docx(title, content)
    
    elif doc_type == "xlsx":
        return _generate_xlsx(title, content)
    
    elif doc_type == "pptx":
        return _generate_pptx(title, content)
    
    else:
        raise ValueError(
            f"Unsupported document type '{doc_type}'. "
            "Supported types: pdf, docx, xlsx, pptx"
        )


